from flask import Flask, render_template, request, jsonify, send_file
import openai
import assemblyai as aai
import os
from dotenv import load_dotenv
from reportlab.pdfgen import canvas
from pptx import Presentation
import io
import requests

# Load environment variables
load_dotenv()

app = Flask(__name__)

# Configure API keys
openai.api_key = os.getenv('OPENAI_API_KEY')
aai.settings.api_key = os.getenv('ASSEMBLYAI_API_KEY')

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/generate_content', methods=['POST'])
def generate_content():
    data = request.json
    text = data.get('text', '')
    platform = data.get('platform', '')
    
    # Generate content using OpenAI
    prompt = f"Create a {platform} post about: {text}. Include relevant hashtags and make it viral."
    response = openai.ChatCompletion.create(
        model="gpt-3.5-turbo",
        messages=[
            {"role": "system", "content": "You are a social media expert."},
            {"role": "user", "content": prompt}
        ]
    )
    
    generated_content = response.choices[0].message.content
    
    return jsonify({'content': generated_content})

@app.route('/generate_image', methods=['POST'])
def generate_image():
    data = request.json
    text = data.get('text', '')
    
    # Generate image using DALL-E
    response = openai.Image.create(
        prompt=text,
        n=1,
        size="512x512"
    )
    
    image_url = response['data'][0]['url']
    return jsonify({'image_url': image_url})

@app.route('/export_pdf', methods=['POST'])
def export_pdf():
    data = request.json
    content = data.get('content', '')
    
    # Create PDF
    buffer = io.BytesIO()
    p = canvas.Canvas(buffer)
    p.drawString(100, 750, content)
    p.save()
    
    buffer.seek(0)
    return send_file(
        buffer,
        as_attachment=True,
        download_name='content.pdf',
        mimetype='application/pdf'
    )

@app.route('/export_ppt', methods=['POST'])
def export_ppt():
    data = request.json
    content = data.get('content', '')
    
    # Create PPT
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "Generated Content"
    
    content_shape = slide.placeholders[1]
    content_shape.text = content
    
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    
    return send_file(
        buffer,
        as_attachment=True,
        download_name='content.pptx',
        mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
    )

@app.route('/tts', methods=['POST'])
def text_to_speech():
    try:
        data = request.json
        text = data.get('text', '')
        language = data.get('language', 'en')

        if not text:
            return jsonify({'error': 'No text provided'}), 400

        # Select voice based on language
        voice = 'alloy'  # Default voice
        if language == 'arabic':
            voice = 'nova'  # Better for Arabic

        # Call OpenAI TTS API
        response = requests.post(
            'https://api.openai.com/v1/audio/speech',
            headers={
                'Authorization': f'Bearer {openai.api_key}',
                'Content-Type': 'application/json'
            },
            json={
                'model': 'tts-1',
                'input': text,
                'voice': voice,
                'response_format': 'mp3'
            }
        )

        if response.status_code != 200:
            error_message = response.json().get('error', {}).get('message', 'Unknown error')
            return jsonify({'error': 'TTS failed', 'details': error_message}), 500

        # Return the audio file
        return send_file(
            io.BytesIO(response.content),
            mimetype='audio/mp3',
            as_attachment=False,
            download_name='output.mp3'
        )

    except Exception as e:
        print(f"Error in TTS: {str(e)}")  # This will show in your server logs
        return jsonify({'error': 'TTS failed', 'details': str(e)}), 500

if __name__ == '__main__':
    app.run(debug=True) 